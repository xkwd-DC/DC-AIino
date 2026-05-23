"""Generate CRAIC defense PPT (18 slides, 16:9)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


REPO = Path(__file__).resolve().parents[1]
FIG = REPO / "docs" / "figures"
OUT = REPO / "docs" / "craic" / "2026-05-23" / "稷韧云图_CRAIC答辩_v1.pptx"

# Palette
DEEP_GREEN = RGBColor(0x1F, 0x4D, 0x2E)        # primary - agriculture
DATA_BLUE = RGBColor(0x1E, 0x5A, 0x8E)         # secondary - technology
ALERT_ORANGE = RGBColor(0xD9, 0x76, 0x06)      # accent - risk / highlight
INK = RGBColor(0x1A, 0x1A, 0x1A)               # body text
SOFT_INK = RGBColor(0x4A, 0x55, 0x52)          # secondary body
PAPER = RGBColor(0xFA, 0xFA, 0xF7)             # slide background
PALE_GREEN = RGBColor(0xE8, 0xF0, 0xE6)        # band / accent block
PALE_BLUE = RGBColor(0xE6, 0xEE, 0xF6)
DIVIDER = RGBColor(0xCF, 0xD6, 0xCC)

ZH_FONT = "Microsoft YaHei"
EN_FONT = "Arial"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- helpers ----------

def set_bg(slide, color=PAPER):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=ZH_FONT,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    tf.word_wrap = True

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
        # ensure east-asian font binding
        rPr = run._r.get_or_add_rPr()
        from lxml import etree
        for tag in ("eastAsia", "cs", "ascii"):
            existing = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
            if existing is not None:
                rPr.remove(existing)
            el = etree.SubElement(
                rPr,
                "{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag,
            )
            el.set("typeface", font if tag != "ascii" else EN_FONT)
    return tb


def add_page_chrome(slide, page_num, total=18, section_label=None):
    # top accent strip
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.08), DEEP_GREEN)
    # footer band
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.42), SLIDE_W, Inches(0.42), DEEP_GREEN)
    add_text(
        slide,
        Inches(0.5),
        SLIDE_H - Inches(0.42),
        Inches(8),
        Inches(0.42),
        "稷韧云图 · 极端气候下粮食生产风险核心因子识别与韧性提升可视化系统",
        size=10,
        color=PAPER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        SLIDE_W - Inches(2.5),
        SLIDE_H - Inches(0.42),
        Inches(2),
        Inches(0.42),
        f"{page_num:02d} / {total:02d}",
        size=10,
        color=PAPER,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if section_label:
        add_text(
            slide,
            Inches(0.5),
            Inches(0.25),
            Inches(6),
            Inches(0.4),
            section_label,
            size=11,
            bold=True,
            color=ALERT_ORANGE,
        )


def add_title(slide, title, subtitle=None, top=Inches(0.7)):
    add_text(
        slide, Inches(0.5), top, Inches(12.3), Inches(0.7),
        title, size=32, bold=True, color=DEEP_GREEN,
    )
    # underline accent bar
    add_rect(slide, Inches(0.5), top + Inches(0.85), Inches(0.6), Inches(0.06), ALERT_ORANGE)
    if subtitle:
        add_text(
            slide, Inches(0.5), top + Inches(0.98), Inches(12.3), Inches(0.4),
            subtitle, size=14, color=SOFT_INK,
        )


def add_bullet(slide, left, top, width, height, items, *, size=16, color=INK,
               bullet_color=ALERT_ORANGE, gap=0.05):
    """items: list of (label, body) or strings."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.3
        p.space_after = Pt(6)
        # bullet glyph
        r0 = p.add_run()
        r0.text = "▸ "
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        r0.font.name = ZH_FONT
        if isinstance(item, tuple):
            label, body = item
            r1 = p.add_run()
            r1.text = label
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = DEEP_GREEN
            r1.font.name = ZH_FONT
            r2 = p.add_run()
            r2.text = "  " + body
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = ZH_FONT
        else:
            r1 = p.add_run()
            r1.text = item
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = ZH_FONT
    return tb


def add_chip(slide, left, top, text, *, fill=PALE_GREEN, text_color=DEEP_GREEN, size=12, width=None):
    if width is None:
        width = Inches(0.06 * len(text) + 0.5)
    height = Inches(0.32)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.adjustments[0] = 0.4
    tf = shp.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = ZH_FONT
    return shp, width


# ---------- slides ----------

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PAPER)
    # left green panel
    add_rect(s, Emu(0), Emu(0), Inches(4.2), SLIDE_H, DEEP_GREEN)
    # diagonal accent
    add_rect(s, Inches(4.2), Emu(0), Inches(0.12), SLIDE_H, ALERT_ORANGE)

    add_text(s, Inches(0.5), Inches(0.5), Inches(3.5), Inches(0.4),
             "CRAIC · 2026", size=14, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(0.5), Inches(0.95), Inches(3.5), Inches(2.5),
             "第二十八届\n中国机器人及\n人工智能创新大赛",
             size=22, bold=True, color=PAPER, line_spacing=1.25)
    add_text(s, Inches(0.5), Inches(3.3), Inches(3.5), Inches(0.4),
             "人工智能创新比赛", size=14, color=PALE_GREEN)

    add_text(s, Inches(0.5), SLIDE_H - Inches(2.2), Inches(3.5), Inches(0.35),
             "团队名称  TEAM", size=10, color=PALE_GREEN)
    add_text(s, Inches(0.5), SLIDE_H - Inches(1.9), Inches(3.5), Inches(0.5),
             "稷韧云图", size=24, bold=True, color=PAPER)
    add_text(s, Inches(0.5), SLIDE_H - Inches(1.3), Inches(3.5), Inches(0.35),
             "潘妙齐 · 熊鑫 · 常宇璇", size=12, color=PALE_GREEN)
    add_text(s, Inches(0.5), SLIDE_H - Inches(0.85), Inches(3.5), Inches(0.35),
             "指导教师：徐苗 / 邱东芳", size=11, color=PALE_GREEN)
    add_text(s, Inches(0.5), SLIDE_H - Inches(0.45), Inches(3.5), Inches(0.35),
             "2026.05", size=10, color=PALE_GREEN)

    # right title area
    add_text(s, Inches(4.8), Inches(1.3), Inches(8.2), Inches(0.5),
             "AI · 多模态 · 可解释 · 决策支持",
             size=14, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(4.8), Inches(1.85), Inches(8.2), Inches(2.6),
             "基于多模态的\n极端气候下我国粮食生产风险的\n核心因子识别与韧性提升\n可视化系统",
             size=30, bold=True, color=DEEP_GREEN, line_spacing=1.3)

    # slogan card
    add_rect(s, Inches(4.8), Inches(5.3), Inches(7.8), Inches(1.0), PALE_GREEN)
    add_rect(s, Inches(4.8), Inches(5.3), Inches(0.1), Inches(1.0), ALERT_ORANGE)
    add_text(s, Inches(5.05), Inches(5.45), Inches(7.4), Inches(0.4),
             "项目愿景", size=11, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(5.05), Inches(5.75), Inches(7.4), Inches(0.5),
             "让 AI 看见每一粒粮食的气候风险",
             size=18, bold=True, color=DEEP_GREEN)


def slide_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 2, section_label="一分钟看懂")
    add_title(s, "一分钟看懂我们做了什么", "Data × Method × Application 的端到端闭环")

    # 3 pillars
    cards = [
        ("01", "数据 DATA",
         "融合 4 类异构数据\n（统计 + 遥感 + 气象 + GIS）\n全国 31 省 · 2011–2023 · 403 样本"),
        ("02", "方法 METHOD",
         "XGBoost-SHAP × Attention-LSTM\n双引擎协同验证\n非线性归因 + 时序响应 双视角"),
        ("03", "应用 APPLICATION",
         "Vue + Flask 可视化决策系统\n四大模块 · 在线交互\n单次推理 ≤ 2 秒"),
    ]
    card_w = Inches(3.9)
    gap = Inches(0.25)
    start_left = Inches(0.6)
    top = Inches(2.3)
    for i, (no, title, body) in enumerate(cards):
        L = start_left + (card_w + gap) * i
        add_rect(s, L, top, card_w, Inches(3.3), PALE_GREEN if i % 2 == 0 else PALE_BLUE)
        add_rect(s, L, top, card_w, Inches(0.08), ALERT_ORANGE)
        add_text(s, L + Inches(0.3), top + Inches(0.3), Inches(3.3), Inches(0.5),
                 no, size=34, bold=True, color=ALERT_ORANGE)
        add_text(s, L + Inches(0.3), top + Inches(1.0), Inches(3.3), Inches(0.4),
                 title, size=16, bold=True, color=DEEP_GREEN)
        add_text(s, L + Inches(0.3), top + Inches(1.5), Inches(3.3), Inches(1.8),
                 body, size=13, color=INK, line_spacing=1.45)

    # bottom strip
    add_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.7), DEEP_GREEN)
    add_text(s, Inches(0.85), Inches(6.05), Inches(11.7), Inches(0.6),
             "破解传统研究「数据单一 · 黑箱模型 · 难以落地」三大瓶颈",
             size=14, bold=True, color=PAPER, anchor=MSO_ANCHOR.MIDDLE)


def slide_background_motivation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 3, section_label="研究背景")
    add_title(s, "为什么做这件事", "政策牵引 × 现实痛点 × 趋势加剧")

    # 3 columns
    items = [
        ("政策牵引",
         "党的二十大、「十四五」规划纲要、历年中央一号文件\n→ 「夯实粮食安全根基,强化产能基础」\n\nIPCC AR6: 气候变化加剧全球农业系统脆弱性",
         DEEP_GREEN, PALE_GREEN),
        ("现实痛点",
         "2023 京津冀「23·7」极端暴雨\n→ 河北 24 万公顷农田受灾\n→ 绝产 13 万公顷, 玉米损失 22 万吨\n\n海河流域 1963 年来首次流域性洪水",
         ALERT_ORANGE, RGBColor(0xFC, 0xEC, 0xD8)),
        ("趋势加剧",
         "2023 全国均温 10.7℃ 创 1961 来新高\n13 省创历史最高温纪录\n\n长江中下游高温干旱、西南「旱涝急转」频发",
         DATA_BLUE, PALE_BLUE),
    ]
    card_w = Inches(3.9)
    gap = Inches(0.25)
    top = Inches(2.2)
    for i, (title, body, header_color, bg) in enumerate(items):
        L = Inches(0.6) + (card_w + gap) * i
        add_rect(s, L, top, card_w, Inches(4.5), bg)
        add_rect(s, L, top, card_w, Inches(0.8), header_color)
        add_text(s, L + Inches(0.3), top + Inches(0.18), card_w, Inches(0.5),
                 title, size=18, bold=True, color=PAPER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, L + Inches(0.3), top + Inches(1.1), card_w - Inches(0.5), Inches(3.3),
                 body, size=12, color=INK, line_spacing=1.5)


def slide_three_gaps(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 4, section_label="研究问题")
    add_title(s, "现有研究的三大瓶颈", "Data Gap · Method Gap · Application Gap")

    rows = [
        ("数据维度单一",
         "依赖单一统计数据\n缺乏遥感、气象、空间矢量等多模态融合",
         "宏观—中观—微观多粒度刻画不足"),
        ("模型可解释性不足",
         "传统机器学习「黑箱化」\n多采用单一模型,缺乏多模型协同验证",
         "因子影响路径与贡献度难以量化"),
        ("研究成果难以落地",
         "多停留在算法与论文层面\n缺乏面向决策者的工程化工具",
         "前沿 AI 难以服务政策实践"),
    ]
    top = Inches(2.2)
    row_h = Inches(1.4)
    for i, (g, desc, impact) in enumerate(rows):
        y = top + row_h * i
        # left label
        add_rect(s, Inches(0.6), y, Inches(2.5), row_h - Inches(0.15), DEEP_GREEN)
        add_text(s, Inches(0.6), y, Inches(2.5), row_h - Inches(0.15),
                 g, size=18, bold=True, color=PAPER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # middle desc
        add_rect(s, Inches(3.2), y, Inches(5.5), row_h - Inches(0.15), PALE_GREEN)
        add_text(s, Inches(3.35), y + Inches(0.18), Inches(5.3), row_h - Inches(0.4),
                 desc, size=13, color=INK, line_spacing=1.4)
        # right impact
        add_rect(s, Inches(8.8), y, Inches(3.9), row_h - Inches(0.15), RGBColor(0xFC, 0xEC, 0xD8))
        add_rect(s, Inches(8.8), y, Inches(0.08), row_h - Inches(0.15), ALERT_ORANGE)
        add_text(s, Inches(9.0), y + Inches(0.18), Inches(3.6), Inches(0.3),
                 "影响", size=10, bold=True, color=ALERT_ORANGE)
        add_text(s, Inches(9.0), y + Inches(0.5), Inches(3.6), row_h - Inches(0.7),
                 impact, size=13, color=INK, line_spacing=1.4)

    add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
             "→ 本项目针对三大瓶颈,从数据层 · 方法层 · 应用层 三个维度系统性破题",
             size=14, bold=True, color=ALERT_ORANGE)


def slide_objectives(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 5, section_label="研究目标")
    add_title(s, "研究目标 与 四大研究内容", "覆盖「风险识别 → 因子归因 → 韧性评估 → 可视化决策」全链条")

    # 4 columns
    contents = [
        ("01", "多模态融合\n数据集构建",
         "统计 + MODIS 遥感 + 气象站点 + GIS 矢量\n线性插值 / Butterworth 滤波 / 时空对齐 / Z-score"),
        ("02", "XGBoost-SHAP\n非线性归因",
         "梯度提升树构建风险拟合模型\nSHAP 三维度: 全局 / 个体 / 交互"),
        ("03", "Attention-LSTM\n时序验证",
         "注意力机制聚焦极端气候敏感时点\n作为双重验证模型独立验证 XGB 结果"),
        ("04", "可视化决策\n支持系统",
         "Vue + Flask 四大模块: 地图 / 看板 / 模拟 / 路径\n面向决策者的可操作工具"),
    ]
    card_w = Inches(2.95)
    gap = Inches(0.15)
    top = Inches(2.2)
    for i, (no, name, body) in enumerate(contents):
        L = Inches(0.6) + (card_w + gap) * i
        add_rect(s, L, top, card_w, Inches(4.5), PAPER, line=DIVIDER)
        add_rect(s, L, top, card_w, Inches(0.08), DATA_BLUE if i % 2 == 0 else DEEP_GREEN)
        add_text(s, L + Inches(0.25), top + Inches(0.3), card_w, Inches(0.6),
                 no, size=28, bold=True, color=ALERT_ORANGE)
        add_text(s, L + Inches(0.25), top + Inches(1.0), card_w - Inches(0.3), Inches(1.0),
                 name, size=16, bold=True, color=DEEP_GREEN, line_spacing=1.25)
        add_text(s, L + Inches(0.25), top + Inches(2.3), card_w - Inches(0.3), Inches(2.0),
                 body, size=12, color=SOFT_INK, line_spacing=1.5)


def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 6, section_label="技术路线")
    add_title(s, "总体技术路线", "多模态数据驱动 + 双引擎协同验证 + 可视化系统集成")

    # 4-stage horizontal pipeline
    stages = [
        ("数据采集", "国家统计局\nCMDC 气象\nMODIS 遥感\nGIS 矢量"),
        ("数据预处理", "缺失值补齐\nButterworth\n时空对齐\nZ-score 标准化"),
        ("双引擎建模", "XGBoost + SHAP\n非线性归因\n+ Att-LSTM\n时序验证"),
        ("系统集成", "Flask API\nVue 前端\n四大功能模块\n推理 ≤ 2s"),
    ]
    n = len(stages)
    box_w = Inches(2.7)
    arrow_w = Inches(0.4)
    total_w = box_w * n + arrow_w * (n - 1)
    start = (SLIDE_W - total_w) / 2
    top = Inches(2.7)
    for i, (name, body) in enumerate(stages):
        L = start + (box_w + arrow_w) * i
        add_rect(s, L, top, box_w, Inches(2.4), DATA_BLUE if i % 2 == 0 else DEEP_GREEN)
        add_text(s, L, top + Inches(0.25), box_w, Inches(0.5),
                 name, size=18, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
        add_text(s, L + Inches(0.25), top + Inches(0.85), box_w - Inches(0.5), Inches(1.5),
                 body, size=12, color=PAPER, align=PP_ALIGN.CENTER, line_spacing=1.5)
        if i < n - 1:
            ar = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                L + box_w + Inches(0.05),
                top + Inches(0.9),
                arrow_w - Inches(0.1),
                Inches(0.6),
            )
            ar.fill.solid()
            ar.fill.fore_color.rgb = ALERT_ORANGE
            ar.line.fill.background()

    # bottom callout
    add_rect(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.1), PALE_GREEN)
    add_rect(s, Inches(0.6), Inches(5.7), Inches(0.1), Inches(1.1), ALERT_ORANGE)
    add_text(s, Inches(0.9), Inches(5.85), Inches(11.7), Inches(0.4),
             "核心闭环", size=11, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(0.9), Inches(6.15), Inches(11.7), Inches(0.6),
             "两个原理完全不同的模型独立建模 → 交叉印证核心因子 → 系统化输出决策建议",
             size=15, bold=True, color=DEEP_GREEN)


def slide_dataset(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 7, section_label="方法 · 数据层")
    add_title(s, "多模态数据集构建", "paper_panel_v3:403 样本 × 27 维  · 全国 31 省 · 2011–2023")

    # left: 4 data sources
    add_text(s, Inches(0.6), Inches(2.1), Inches(6), Inches(0.4),
             "数据来源", size=14, bold=True, color=DEEP_GREEN)
    sources = [
        ("农业生产", "粮食单产 / 农机总动力 / 化肥施用 / 灌溉面积", "中国统计年鉴"),
        ("气象监测", "温度 / 降水 / 日照 / SPEI 干旱指数", "CMDC 中国气象数据网"),
        ("致灾因子", "洪涝占比 / 旱灾面积 / 水灾面积", "中国水旱灾害公报"),
        ("遥感空间", "MODIS NDVI / LST / 省级行政矢量 / 耕地分布", "NASA EARTHDATA + 国家地球数据中心"),
    ]
    top = Inches(2.55)
    for i, (cat, vars_, src) in enumerate(sources):
        y = top + Inches(0.9) * i
        add_rect(s, Inches(0.6), y, Inches(0.12), Inches(0.75), ALERT_ORANGE if i == 3 else DATA_BLUE)
        add_text(s, Inches(0.85), y, Inches(1.6), Inches(0.4),
                 cat, size=14, bold=True, color=DEEP_GREEN)
        add_text(s, Inches(0.85), y + Inches(0.35), Inches(5.5), Inches(0.4),
                 vars_, size=11, color=INK)
        add_text(s, Inches(0.85), y + Inches(0.6), Inches(5.5), Inches(0.3),
                 "来源:" + src, size=10, color=SOFT_INK)

    # right: preprocessing pipeline
    add_rect(s, Inches(7.3), Inches(2.1), Inches(5.4), Inches(4.7), PALE_BLUE)
    add_text(s, Inches(7.5), Inches(2.25), Inches(5), Inches(0.4),
             "预处理流程", size=14, bold=True, color=DEEP_GREEN)
    steps = [
        ("① 缺失值补齐", "线性插值,确保时空序列连续性"),
        ("② 单产去趋势", "Butterworth 低通滤波分离技术进步\n提取气候波动的随机残差作为风险 Y"),
        ("③ 时空对齐", "统一年度 + 省级口径\n栅格遥感 → 省域空间聚合"),
        ("④ Z-score 标准化", "消除量纲差异\nμ=0, σ=1"),
    ]
    top2 = Inches(2.7)
    for i, (k, v) in enumerate(steps):
        y = top2 + Inches(1.0) * i
        add_text(s, Inches(7.5), y, Inches(5), Inches(0.35),
                 k, size=12, bold=True, color=ALERT_ORANGE)
        add_text(s, Inches(7.5), y + Inches(0.35), Inches(5), Inches(0.65),
                 v, size=11, color=INK, line_spacing=1.4)

    # bottom hook
    add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
             "→ 致灾维度 × 韧性维度 双维度指标体系,从风险诱发源与对冲力两条逻辑线并行刻画",
             size=12, bold=True, color=ALERT_ORANGE)


def slide_xgb_shap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 8, section_label="方法 · 引擎一")
    add_title(s, "XGBoost-SHAP 非线性归因引擎", "梯度提升树 × Shapley 值,把「黑箱」变「白箱」")

    # left: XGB
    add_rect(s, Inches(0.6), Inches(2.1), Inches(6), Inches(4.7), PAPER, line=DIVIDER)
    add_rect(s, Inches(0.6), Inches(2.1), Inches(6), Inches(0.08), DATA_BLUE)
    add_text(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(0.4),
             "XGBoost 梯度提升树", size=16, bold=True, color=DEEP_GREEN)
    add_text(s, Inches(0.85), Inches(2.75), Inches(5.5), Inches(0.4),
             "处理非线性 · 多耦合特征的集成学习算法",
             size=11, color=SOFT_INK)
    # formula block
    add_rect(s, Inches(0.85), Inches(3.3), Inches(5.4), Inches(1.4), PALE_BLUE)
    add_text(s, Inches(1.0), Inches(3.4), Inches(5.2), Inches(0.4),
             "加法模型", size=10, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(1.0), Inches(3.7), Inches(5.2), Inches(0.5),
             "ŷᵢ = Σ fₖ(xᵢ),  k=1,…,K",
             size=16, bold=True, color=INK, font="Arial")
    add_text(s, Inches(1.0), Inches(4.2), Inches(5.2), Inches(0.4),
             "目标函数(损失 + 正则)",
             size=10, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(1.0), Inches(4.45), Inches(5.2), Inches(0.4),
             "Obj = Σ L(yᵢ, ŷᵢ) + Σ Ω(fₖ)",
             size=14, bold=True, color=INK, font="Arial")
    add_bullet(s, Inches(0.85), Inches(4.95), Inches(5.4), Inches(1.8),
               ["串行构建多棵决策树,优化模型残差",
                "正则项惩罚树复杂度,避免过拟合",
                "对粮食风险这类非线性 · 多耦合数据效果显著"],
               size=11)

    # right: SHAP
    add_rect(s, Inches(6.85), Inches(2.1), Inches(5.85), Inches(4.7), PAPER, line=DIVIDER)
    add_rect(s, Inches(6.85), Inches(2.1), Inches(5.85), Inches(0.08), ALERT_ORANGE)
    add_text(s, Inches(7.1), Inches(2.3), Inches(5.5), Inches(0.4),
             "SHAP 可解释性归因", size=16, bold=True, color=DEEP_GREEN)
    add_text(s, Inches(7.1), Inches(2.75), Inches(5.5), Inches(0.4),
             "基于 Shapley 值的统一归因框架",
             size=11, color=SOFT_INK)

    sh_items = [
        ("全局重要性", "Mean|SHAP value| 对特征排序\n→ 识别主导风险因子"),
        ("个体决策", "瀑布图分解单个样本的预测\n→ 每个省份专属归因报告"),
        ("特征交互", "SHAP 交互依赖图\n→ 揭示因子间非线性耦合"),
    ]
    top = Inches(3.3)
    for i, (k, v) in enumerate(sh_items):
        y = top + Inches(1.1) * i
        add_rect(s, Inches(7.1), y, Inches(0.08), Inches(0.9), ALERT_ORANGE)
        add_text(s, Inches(7.3), y, Inches(5.2), Inches(0.4),
                 k, size=13, bold=True, color=DEEP_GREEN)
        add_text(s, Inches(7.3), y + Inches(0.4), Inches(5.2), Inches(0.7),
                 v, size=11, color=INK, line_spacing=1.4)


def slide_att_lstm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 9, section_label="方法 · 引擎二")
    add_title(s, "Attention-LSTM 时序验证引擎", "LSTM 门控 + 注意力加权,自适应聚焦极端气候敏感时点")

    # left: arch description
    add_rect(s, Inches(0.6), Inches(2.1), Inches(7.5), Inches(4.7), PALE_GREEN)
    add_text(s, Inches(0.85), Inches(2.3), Inches(7), Inches(0.4),
             "架构与机制", size=15, bold=True, color=DEEP_GREEN)
    add_bullet(s, Inches(0.85), Inches(2.75), Inches(7.1), Inches(3.5),
               [
                   ("LSTM 门控记忆", "输入门 / 遗忘门 / 输出门,捕捉长时序依赖"),
                   ("Attention 加权", "softmax 归一(sum=1),自适应聚焦关键时点"),
                   ("时序敏感性", "对洪涝 / 干旱发生月的注意力权重显著放大"),
                   ("双重验证", "作为 XGBoost-SHAP 的独立交叉验证模型"),
               ], size=13)

    # right: highlight card with real prediction
    add_rect(s, Inches(8.4), Inches(2.1), Inches(4.3), Inches(4.7), DEEP_GREEN)
    add_rect(s, Inches(8.4), Inches(2.1), Inches(4.3), Inches(0.6), ALERT_ORANGE)
    add_text(s, Inches(8.4), Inches(2.1), Inches(4.3), Inches(0.6),
             "单样本实测验证", size=14, bold=True, color=PAPER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.6), Inches(2.95), Inches(3.9), Inches(0.4),
             "Case · 河南省 2022", size=12, color=PALE_GREEN)
    add_text(s, Inches(8.6), Inches(3.4), Inches(3.9), Inches(0.4),
             "真实单产", size=11, color=PALE_GREEN)
    add_text(s, Inches(8.6), Inches(3.7), Inches(3.9), Inches(0.6),
             "4615 kg/ha", size=22, bold=True, color=PAPER, font="Arial")
    add_text(s, Inches(8.6), Inches(4.4), Inches(3.9), Inches(0.4),
             "Att-LSTM 预测", size=11, color=PALE_GREEN)
    add_text(s, Inches(8.6), Inches(4.7), Inches(3.9), Inches(0.6),
             "4657.6 kg/ha", size=22, bold=True, color=ALERT_ORANGE, font="Arial")
    # delta
    add_rect(s, Inches(8.6), Inches(5.5), Inches(3.9), Inches(0.9), PALE_GREEN)
    add_text(s, Inches(8.6), Inches(5.55), Inches(3.9), Inches(0.4),
             "预测误差", size=10, bold=True, color=ALERT_ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.6), Inches(5.85), Inches(3.9), Inches(0.5),
             "+42.6 kg/ha", size=18, bold=True, color=DEEP_GREEN,
             align=PP_ALIGN.CENTER, font="Arial")


def slide_dual_engine(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 10, section_label="方法 · 核心创新")
    add_title(s, "双引擎协同验证机制", "两个原理完全不同的模型 → 独立建模 → 交叉印证 → 结论稳健")

    # center diagram: two engines feeding into a converge box
    # left engine
    add_rect(s, Inches(0.6), Inches(2.4), Inches(4.2), Inches(2.6), DATA_BLUE)
    add_text(s, Inches(0.6), Inches(2.55), Inches(4.2), Inches(0.5),
             "Engine A", size=11, bold=True, color=PALE_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.85), Inches(4.2), Inches(0.5),
             "XGBoost + SHAP", size=20, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.85), Inches(3.6), Inches(3.7), Inches(1.3),
             "数据驱动归因\n非线性树模型\n全局/个体/交互三维归因\n生物量主导视角",
             size=12, color=PALE_GREEN, align=PP_ALIGN.CENTER, line_spacing=1.5)

    # right engine
    add_rect(s, Inches(8.5), Inches(2.4), Inches(4.2), Inches(2.6), DEEP_GREEN)
    add_text(s, Inches(8.5), Inches(2.55), Inches(4.2), Inches(0.5),
             "Engine B", size=11, bold=True, color=PALE_GREEN, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.5), Inches(2.85), Inches(4.2), Inches(0.5),
             "Attention-LSTM", size=20, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.75), Inches(3.6), Inches(3.7), Inches(1.3),
             "时序响应归因\n深度时序模型\n注意力加权 softmax\n灾害响应/政策杠杆视角",
             size=12, color=PALE_GREEN, align=PP_ALIGN.CENTER, line_spacing=1.5)

    # arrows
    a1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.85), Inches(3.5), Inches(0.7), Inches(0.4))
    a1.fill.solid(); a1.fill.fore_color.rgb = ALERT_ORANGE; a1.line.fill.background()
    a2 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.75), Inches(3.5), Inches(0.7), Inches(0.4))
    a2.fill.solid(); a2.fill.fore_color.rgb = ALERT_ORANGE; a2.line.fill.background()

    # central converge
    add_rect(s, Inches(5.55), Inches(2.6), Inches(2.25), Inches(2.2), ALERT_ORANGE)
    add_text(s, Inches(5.55), Inches(2.8), Inches(2.25), Inches(0.5),
             "交叉印证", size=11, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.55), Inches(3.1), Inches(2.25), Inches(0.6),
             "稳健归因", size=20, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.55), Inches(3.8), Inches(2.25), Inches(1.0),
             "Top-K 因子\nRobust to\nmodel choice",
             size=10, color=PAPER, align=PP_ALIGN.CENTER, line_spacing=1.4, font="Arial")

    # bottom narrative
    add_rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.4), PALE_GREEN)
    add_text(s, Inches(0.85), Inches(5.55), Inches(11.7), Inches(0.4),
             "关键叙事", size=11, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(0.85), Inches(5.85), Inches(11.7), Inches(0.95),
             "两个完全不同原理的模型相互印证 → 既看「数据生长状态」也看「灾害政策响应」\n"
             "→ 「互补而非简单一致」,决策视角更全面、结论更稳健",
             size=15, bold=True, color=DEEP_GREEN, line_spacing=1.5)


def slide_results_r2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 11, section_label="实验结果 · 精度")
    add_title(s, "三模型精度对比", "10 种子稳健性实验 · 全部远超申报书硬指标")

    # table-like grid
    headers = ["模型", "Test R² (10 seeds 均值 ± std)", "RMSE (kg/ha)", "硬指标达成"]
    rows = [
        ("XGBoost", "0.9072 ± 0.0383", "312.5", "≥ 0.62  远超 ✅"),
        ("LSTM", "0.8856 ± 0.0223", "362.0", "≥ 0.62  远超 ✅"),
        ("Attention-LSTM", "0.8160 ± 0.0502", "456.4", "≥ 0.65  通过 ✅"),
    ]
    table_left = Inches(0.8)
    table_top = Inches(2.3)
    col_w = [Inches(2.6), Inches(4.0), Inches(2.4), Inches(3.0)]
    row_h = Inches(0.7)

    # header
    L = table_left
    for i, h in enumerate(headers):
        add_rect(s, L, table_top, col_w[i], row_h, DEEP_GREEN)
        add_text(s, L, table_top, col_w[i], row_h,
                 h, size=12, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        L += col_w[i]

    # body
    for r, row in enumerate(rows):
        L = table_left
        y = table_top + row_h * (r + 1)
        # alternating fill, highlight Attention-LSTM
        fill = ALERT_ORANGE if r == 2 else (PALE_GREEN if r % 2 == 0 else PAPER)
        text_color = PAPER if r == 2 else INK
        for i, cell in enumerate(row):
            add_rect(s, L, y, col_w[i], row_h, fill, line=DIVIDER)
            bold = (i == 0 or i == 1)
            add_text(s, L, y, col_w[i], row_h, cell,
                     size=14 if i in (1,) else 13,
                     bold=bold, color=text_color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font="Arial" if i in (1, 2) else ZH_FONT)
            L += col_w[i]

    # callout below
    add_rect(s, Inches(0.8), Inches(5.6), Inches(12.0), Inches(1.2), PALE_BLUE)
    add_text(s, Inches(1.05), Inches(5.7), Inches(11.5), Inches(0.4),
             "亮点", size=11, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(1.05), Inches(6.0), Inches(11.5), Inches(0.8),
             "XGBoost / LSTM / Attention-LSTM 三模型 R² 全部 ≥ 0.81;\n"
             "Att-LSTM 在河南 2022 单样本误差仅 +42.6 kg/ha,模型已具备落地预测能力",
             size=13, color=INK, line_spacing=1.45)


def slide_factor_complement(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 12, section_label="实验结果 · 归因")
    add_title(s, "核心因子识别:互补而非简单一致",
              "双引擎 Top-3 仅 Temp 重合 → 数据驱动 × 灾害响应 双视角")

    # left: XGB SHAP Top-3
    add_rect(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.7), PAPER, line=DIVIDER)
    add_rect(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(0.6), DATA_BLUE)
    add_text(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(0.6),
             "Engine A · XGBoost SHAP Top-3",
             size=14, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    xgb_items = [
        ("#1  NDVI",      "植被指数异常 · 生物量主导"),
        ("#2  Prec",      "年累计降水量"),
        ("#3  Temp",      "年平均气温(★ 与 B 重合)"),
    ]
    for i, (k, v) in enumerate(xgb_items):
        y = Inches(2.95) + Inches(1.2) * i
        add_rect(s, Inches(0.85), y, Inches(0.4), Inches(1.0),
                 ALERT_ORANGE if i == 2 else DATA_BLUE)
        add_text(s, Inches(0.85), y, Inches(0.4), Inches(1.0),
                 f"{i+1}", size=22, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Arial")
        add_text(s, Inches(1.4), y, Inches(5.0), Inches(0.5),
                 k.split("  ")[1], size=18, bold=True,
                 color=ALERT_ORANGE if i == 2 else DEEP_GREEN)
        add_text(s, Inches(1.4), y + Inches(0.5), Inches(5.0), Inches(0.5),
                 v, size=11, color=SOFT_INK)

    add_text(s, Inches(0.85), Inches(6.45), Inches(5.6), Inches(0.3),
             "视角:数据驱动 · 生物量主导", size=11, bold=True, color=DATA_BLUE)

    # right: Att-LSTM Top-3
    add_rect(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(4.7), PAPER, line=DIVIDER)
    add_rect(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(0.6), DEEP_GREEN)
    add_text(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(0.6),
             "Engine B · Att-LSTM Attention Top-3",
             size=14, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    att_items = [
        ("#1  Drou_A",  "干旱面积 · 灾害响应主导"),
        ("#2  SPEI",    "标准化降水蒸散指数"),
        ("#3  Temp",    "年平均气温(★ 与 A 重合)"),
    ]
    for i, (k, v) in enumerate(att_items):
        y = Inches(2.95) + Inches(1.2) * i
        add_rect(s, Inches(7.1), y, Inches(0.4), Inches(1.0),
                 ALERT_ORANGE if i == 2 else DEEP_GREEN)
        add_text(s, Inches(7.1), y, Inches(0.4), Inches(1.0),
                 f"{i+1}", size=22, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Arial")
        add_text(s, Inches(7.65), y, Inches(5.0), Inches(0.5),
                 k.split("  ")[1], size=18, bold=True,
                 color=ALERT_ORANGE if i == 2 else DEEP_GREEN)
        add_text(s, Inches(7.65), y + Inches(0.5), Inches(5.0), Inches(0.5),
                 v, size=11, color=SOFT_INK)

    add_text(s, Inches(7.1), Inches(6.45), Inches(5.6), Inches(0.3),
             "视角:灾害响应 · policy lever", size=11, bold=True, color=DEEP_GREEN)


def slide_interaction_ablation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 13, section_label="实验结果 · 深度归因")
    add_title(s, "SHAP 交互效应 与 多模态融合增益",
              "灌溉×气温非线性耦合 + 数据融合带来 4.9–6.7 pp R² 增益")

    # left figure (SHAP interaction)
    fig = FIG / "fig6_shap_interaction.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(0.5), Inches(2.15),
                             width=Inches(6.0))
    add_text(s, Inches(0.5), Inches(5.65), Inches(6.0), Inches(0.4),
             "图 · 灌溉率 × 气温 SHAP 交互依赖", size=11, bold=True, color=DEEP_GREEN)
    add_text(s, Inches(0.5), Inches(5.95), Inches(6.0), Inches(0.85),
             "→ 灌溉提升整体降低风险;但在高温环境下灌溉减损效果被部分抵消\n"
             "→ 启示:单一提升灌溉率不够,需结合温度调节综合施策",
             size=10, color=INK, line_spacing=1.45)

    # right: ablation table
    add_text(s, Inches(7.0), Inches(2.15), Inches(6), Inches(0.4),
             "多模态融合消融实验 (XGBoost R²)",
             size=14, bold=True, color=DEEP_GREEN)

    rows = [
        ("仅统计数据", "0.498", PALE_BLUE),
        ("+ 遥感", "0.522", PALE_BLUE),
        ("+ 遥感 + 气象", "0.547", PALE_BLUE),
        ("全模态融合", "0.5647", ALERT_ORANGE),
    ]
    top = Inches(2.7)
    for i, (k, v, c) in enumerate(rows):
        y = top + Inches(0.75) * i
        text_color = PAPER if c == ALERT_ORANGE else INK
        add_rect(s, Inches(7.0), y, Inches(3.8), Inches(0.6), c)
        add_text(s, Inches(7.2), y, Inches(3.6), Inches(0.6), k,
                 size=13, bold=(i == 3), color=text_color, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(10.8), y, Inches(2.0), Inches(0.6), c, line=DIVIDER)
        add_text(s, Inches(10.8), y, Inches(2.0), Inches(0.6), v,
                 size=15, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Arial")

    # callout
    add_rect(s, Inches(7.0), Inches(5.85), Inches(5.8), Inches(0.95), DEEP_GREEN)
    add_text(s, Inches(7.2), Inches(5.95), Inches(5.4), Inches(0.4),
             "结论", size=10, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(7.2), Inches(6.2), Inches(5.4), Inches(0.6),
             "多模态融合带来 4.9–6.7 pp R² 提升,\n"
             "证明数据融合的实质性价值",
             size=12, bold=True, color=PAPER, line_spacing=1.4)


def slide_system(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 14, section_label="成果落地")
    add_title(s, "可视化决策系统 · 四大模块",
              "Vue 3 + Flask + PostgreSQL · 单次推理 ≤ 2 秒")

    modules = [
        ("M1", "风险时空地图",
         "ECharts 中国地图 + 时间滑块\n查询任意年份省域风险分布 + Top10 排行",
         DATA_BLUE),
        ("M2", "SHAP 归因看板",
         "瀑布图 + 蜂群图 + 特征重要性\n点击省份查看专属个体归因报告",
         DEEP_GREEN),
        ("M3", "参数情景模拟",
         "反事实推理:灌溉率 +10% / 化肥 -20%\n实时计算风险变化曲线",
         ALERT_ORANGE),
        ("M4", "韧性路径推荐",
         "11 条规则引擎 + 三阶段时间线\n生成差异化政策建议(短/中/长期)",
         RGBColor(0x70, 0x4C, 0x9F)),
    ]
    card_w = Inches(5.9)
    card_h = Inches(2.15)
    gap = Inches(0.25)
    top = Inches(2.1)
    for idx, (no, name, body, color) in enumerate(modules):
        r, c = divmod(idx, 2)
        L = Inches(0.6) + (card_w + gap) * c
        T = top + (card_h + Inches(0.2)) * r
        add_rect(s, L, T, card_w, card_h, PAPER, line=DIVIDER)
        add_rect(s, L, T, Inches(0.12), card_h, color)
        add_text(s, L + Inches(0.35), T + Inches(0.25), Inches(0.8), Inches(0.4),
                 no, size=11, bold=True, color=ALERT_ORANGE)
        add_text(s, L + Inches(0.35), T + Inches(0.55), Inches(5.4), Inches(0.5),
                 name, size=18, bold=True, color=DEEP_GREEN)
        add_text(s, L + Inches(0.35), T + Inches(1.1), Inches(5.4), Inches(1.0),
                 body, size=12, color=INK, line_spacing=1.5)

    # bottom perf bar
    add_rect(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.55), DEEP_GREEN)
    add_text(s, Inches(0.85), Inches(6.75), Inches(11.7), Inches(0.55),
             "系统性能:推理响应 ≤ 2s   |   全国 31 省覆盖   |   开放 API   |   面向决策者的可操作工具",
             size=12, bold=True, color=PAPER, anchor=MSO_ANCHOR.MIDDLE)


def slide_innovations(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 15, section_label="创新点总结")
    add_title(s, "五大创新点", "数据层 · 方法层 × 2 · 应用层 · 贡献层")

    items = [
        ("01", "数据层",
         "多模态异构数据融合的全国尺度风险识别框架",
         "首次构建覆盖统计 + MODIS 遥感 + 气象 + GIS 四类异构数据源的融合体系",
         DATA_BLUE),
        ("02", "方法层 ①",
         "注意力增强的 XGBoost-SHAP × LSTM 双引擎协同验证",
         "首次将可解释机器学习与注意力深度学习时序模型相结合,破解黑箱困境",
         DEEP_GREEN),
        ("03", "方法层 ②",
         "SHAP 三维度归因 + 多维交互效应深度挖掘",
         "全局重要性 / 个体决策 / 特征交互三维系统性解构,揭示复杂耦合机理",
         ALERT_ORANGE),
        ("04", "应用层",
         "全流程可视化决策支持系统的工程化落地",
         "首次将「多模态识别 + 可解释 AI + 决策支持」三位一体集成到单一平台",
         RGBColor(0x70, 0x4C, 0x9F)),
        ("05", "贡献层",
         "多模态数据集构建与处理流程的规范化输出",
         "制定可复用的数据基础与方法参考,服务后续相关研究",
         RGBColor(0x0E, 0x7C, 0x86)),
    ]
    top = Inches(2.2)
    row_h = Inches(0.85)
    for i, (no, layer, title, body, color) in enumerate(items):
        y = top + (row_h + Inches(0.05)) * i
        add_rect(s, Inches(0.6), y, Inches(0.95), row_h, color)
        add_text(s, Inches(0.6), y, Inches(0.95), row_h, no,
                 size=22, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Arial")
        add_rect(s, Inches(1.55), y, Inches(2.1), row_h, PALE_GREEN)
        add_text(s, Inches(1.55), y, Inches(2.1), row_h, layer,
                 size=13, bold=True, color=DEEP_GREEN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.85), y + Inches(0.08), Inches(8.9), Inches(0.4),
                 title, size=14, bold=True, color=DEEP_GREEN)
        add_text(s, Inches(3.85), y + Inches(0.45), Inches(8.9), Inches(0.4),
                 body, size=11, color=SOFT_INK, line_spacing=1.4)


def slide_value(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 16, section_label="应用前景")
    add_title(s, "应用前景 与 社会价值",
              "服务国家粮食安全战略 · 推动可解释 AI 在农业领域落地")

    items = [
        ("🎯", "服务国家粮食安全战略",
         "为政府农业管理部门提供风险预警与资源调配的科学依据\n落实党的二十大「夯实粮食安全根基」部署",
         DEEP_GREEN),
        ("💼", "赋能农业保险与产业",
         "支持保险产品差异化设计、农资精准营销、费率精算与赔付决策",
         DATA_BLUE),
        ("🤖", "推动可解释 AI 在农业落地",
         "全国尺度「多模态识别 + 可解释 AI + 决策支持」三位一体示范样板",
         ALERT_ORANGE),
        ("🗺️", "支撑区域差异化政策",
         "针对北方旱区、长江洪涝区、西南「旱涝急转」区提出差异化韧性路径",
         RGBColor(0x70, 0x4C, 0x9F)),
        ("🔓", "促进研究开放协作",
         "数据集与系统 API 开放,推动跨学科、跨区域对比研究",
         RGBColor(0x0E, 0x7C, 0x86)),
    ]
    top = Inches(2.2)
    h = Inches(0.86)
    for i, (emoji, title, body, color) in enumerate(items):
        y = top + (h + Inches(0.06)) * i
        add_rect(s, Inches(0.6), y, Inches(0.95), h, color)
        add_text(s, Inches(0.6), y, Inches(0.95), h, emoji,
                 size=28, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Arial")
        add_rect(s, Inches(1.55), y, Inches(11.2), h, PAPER, line=DIVIDER)
        add_text(s, Inches(1.85), y + Inches(0.1), Inches(10.7), Inches(0.4),
                 title, size=15, bold=True, color=DEEP_GREEN)
        add_text(s, Inches(1.85), y + Inches(0.45), Inches(10.7), Inches(0.45),
                 body, size=11, color=SOFT_INK, line_spacing=1.4)


def slide_future(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_page_chrome(s, 17, section_label="改进方向")
    add_title(s, "存在问题 与 今后改进方向",
              "模型 / 数据 / 系统 / 应用 四维度持续演进")

    items = [
        ("模型层面", "罕见极端事件识别能力有限",
         "引入 Transformer + 多模型 Stacking\n贝叶斯神经网络给出不确定性区间",
         DATA_BLUE),
        ("数据层面", "省级粒度偏粗,西部站点稀疏",
         "下沉至地市级 / 县域级\n引入 Sentinel-2 + 高分系列高分辨率遥感",
         DEEP_GREEN),
        ("系统层面", "移动端适配弱,反事实推理基于静态规则",
         "响应式 + 微信小程序\nDo-Calculus / 合成控制法因果推断增强",
         ALERT_ORANGE),
        ("应用层面", "韧性路径与政策工具对接不够精细",
         "与地方农业局、保险机构试点合作\n开发面向农户的简化版客户端",
         RGBColor(0x70, 0x4C, 0x9F)),
    ]
    card_w = Inches(5.95)
    card_h = Inches(2.3)
    gap = Inches(0.2)
    top = Inches(2.15)
    for idx, (layer, problem, plan, color) in enumerate(items):
        r, c = divmod(idx, 2)
        L = Inches(0.6) + (card_w + gap) * c
        T = top + (card_h + Inches(0.15)) * r
        add_rect(s, L, T, card_w, card_h, PAPER, line=DIVIDER)
        add_rect(s, L, T, card_w, Inches(0.5), color)
        add_text(s, L + Inches(0.25), T, card_w, Inches(0.5), layer,
                 size=14, bold=True, color=PAPER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, L + Inches(0.25), T + Inches(0.65), card_w - Inches(0.5), Inches(0.4),
                 "❗ 问题", size=10, bold=True, color=ALERT_ORANGE)
        add_text(s, L + Inches(0.25), T + Inches(0.95), card_w - Inches(0.5), Inches(0.5),
                 problem, size=12, color=INK)
        add_text(s, L + Inches(0.25), T + Inches(1.45), card_w - Inches(0.5), Inches(0.4),
                 "✓ 改进方向", size=10, bold=True, color=DEEP_GREEN)
        add_text(s, L + Inches(0.25), T + Inches(1.75), card_w - Inches(0.5), Inches(0.5),
                 plan, size=12, color=SOFT_INK, line_spacing=1.4)


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DEEP_GREEN)
    # accent strip
    add_rect(s, Emu(0), Inches(1.0), SLIDE_W, Inches(0.06), ALERT_ORANGE)

    add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5),
             "Thanks · Q&A", size=14, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(1.0),
             "总结与致谢",
             size=44, bold=True, color=PAPER)

    summary = [
        ("数据",
         "构建覆盖全国 31 省 13 年的多模态融合数据集 paper_panel_v3"),
        ("方法",
         "建立 XGBoost-SHAP × Attention-LSTM 双引擎协同验证 + SHAP 三维归因"),
        ("落地",
         "开发端到端可视化决策系统,四大模块、推理 ≤ 2s、推动可解释 AI 在农业落地"),
    ]
    top = Inches(2.6)
    for i, (k, v) in enumerate(summary):
        y = top + Inches(0.85) * i
        add_rect(s, Inches(0.6), y, Inches(1.3), Inches(0.65), ALERT_ORANGE)
        add_text(s, Inches(0.6), y, Inches(1.3), Inches(0.65), k,
                 size=15, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.05), y, Inches(11.0), Inches(0.65), v,
                 size=14, color=PAPER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

    # acknowledgments
    add_rect(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.3), RGBColor(0x14, 0x36, 0x21))
    add_text(s, Inches(0.85), Inches(5.65), Inches(11.7), Inches(0.4),
             "致谢", size=11, bold=True, color=ALERT_ORANGE)
    add_text(s, Inches(0.85), Inches(5.95), Inches(11.7), Inches(0.45),
             "感谢徐苗老师、邱东芳老师的悉心指导;感谢学院与团队各位成员的协作支持",
             size=13, color=PAPER)
    add_text(s, Inches(0.85), Inches(6.4), Inches(11.7), Inches(0.4),
             "团队 · 稷韧云图   |   潘妙齐 · 熊鑫 · 常宇璇   |   敬请评委批评指正",
             size=12, bold=True, color=PALE_GREEN)


# ---------- main ----------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)                  # 01
    slide_overview(prs)               # 02
    slide_background_motivation(prs)  # 03
    slide_three_gaps(prs)             # 04
    slide_objectives(prs)             # 05
    slide_pipeline(prs)               # 06
    slide_dataset(prs)                # 07
    slide_xgb_shap(prs)               # 08
    slide_att_lstm(prs)               # 09
    slide_dual_engine(prs)            # 10
    slide_results_r2(prs)             # 11
    slide_factor_complement(prs)      # 12
    slide_interaction_ablation(prs)   # 13
    slide_system(prs)                 # 14
    slide_innovations(prs)            # 15
    slide_value(prs)                  # 16
    slide_future(prs)                 # 17
    slide_closing(prs)                # 18

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"OK · saved to {OUT}")
    print(f"size: {OUT.stat().st_size/1024:.1f} KB · slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
